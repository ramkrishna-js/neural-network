import fitz  # PyMuPDF
from pptx import Presentation
import os
from rich import print as rprint

class NoteProcessor:
    """Extracts text from PDF, PPTX, and TXT files."""
    
    @staticmethod
    def extract_text(file_path):
        ext = os.path.splitext(file_path)[1].lower()
        text = ""
        
        try:
            if ext == '.pdf':
                doc = fitz.open(file_path)
                for page in doc:
                    text += page.get_text()
            elif ext == '.pptx':
                prs = Presentation(file_path)
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            text += shape.text + " "
            elif ext == '.txt':
                with open(file_path, 'r', encoding='utf-8') as f:
                    text = f.read()
            else:
                rprint(f"[bold red]Unsupported file format: {ext}[/bold red]")
                return None
            
            return text
        except Exception as e:
            rprint(f"[bold red]Error processing {file_path}: {e}[/bold red]")
            return None

class Tokenizer:
    """Converts characters to integers and vice versa."""
    def __init__(self, text):
        self.chars = sorted(list(set(text)))
        self.vocab_size = len(self.chars)
        self.char_to_ix = {ch: i for i, ch in enumerate(self.chars)}
        self.ix_to_char = {i: ch for i, ch in enumerate(self.chars)}

    def encode(self, text):
        return [self.char_to_ix[ch] for ch in text if ch in self.char_to_ix]

    def decode(self, rix):
        return "".join([self.ix_to_char[i] for i in rix])
